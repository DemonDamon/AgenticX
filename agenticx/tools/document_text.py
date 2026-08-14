#!/usr/bin/env python3
"""Shared document text extraction for chat and knowledge base.

Author: Damon Li
"""

from __future__ import annotations

import asyncio
import platform
import shutil
from pathlib import Path
from typing import Callable

from agenticx.tools.adapters.liteparse import LiteParseAdapter


NATIVE_READER_EXTS = {".pdf", ".docx", ".pptx"}
LITEPARSE_REQUIRED_EXTS = {
    ".doc",
    ".ppt",
    ".xls",
    ".xlsx",
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".bmp",
}
LIBREOFFICE_REQUIRED_EXTS = {".doc", ".ppt", ".xls", ".xlsx"}
PLAIN_TEXT_EXTS = {
    ".csv",
    ".htm",
    ".html",
    ".js",
    ".json",
    ".log",
    ".markdown",
    ".md",
    ".py",
    ".rst",
    ".ts",
    ".tsv",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}
LITEPARSE_INSTALL_HINT = "npm i -g @llamaindex/liteparse"


class DocumentTextError(RuntimeError):
    """Structured document parse failure with a stable status code."""

    def __init__(
        self,
        code: str,
        user_message: str,
        *,
        install_hint: str | None = None,
    ) -> None:
        super().__init__(user_message)
        self.code = code
        self.user_message = user_message
        self.install_hint = install_hint


def libreoffice_install_hint() -> str:
    """Return platform-specific LibreOffice install command."""
    system = platform.system().strip().lower()
    if system == "darwin":
        return "brew install --cask libreoffice"
    if system == "windows":
        return "choco install libreoffice-fresh"
    return "apt-get install libreoffice"


def libreoffice_available() -> bool:
    """Return True when soffice/libreoffice is on PATH."""
    return bool(shutil.which("soffice") or shutil.which("libreoffice"))


def _read_pdf_text(path: Path) -> str:
    """Read a PDF without importing the local knowledge subsystem."""

    try:
        import fitz  # type: ignore
    except ImportError:
        fitz = None

    if fitz is not None:
        document = fitz.open(str(path))
        try:
            if document.needs_pass:
                raise ValueError("PDF requires password")
            pages = []
            for page_number, page in enumerate(document, 1):
                text = page.get_text()
                if isinstance(text, str) and text.strip():
                    pages.append(f"--- Page {page_number} ---\n{text}")
            return "\n\n".join(pages)
        finally:
            document.close()

    try:
        import pypdf  # type: ignore
    except ImportError as exc:
        raise ImportError(
            "No PDF library available. Install PyMuPDF or pypdf."
        ) from exc

    with path.open("rb") as handle:
        reader = pypdf.PdfReader(handle)
        if reader.is_encrypted:
            raise ValueError("PDF requires password")
        pages = []
        for page_number, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if isinstance(text, str) and text.strip():
                pages.append(f"--- Page {page_number} ---\n{text}")
        return "\n\n".join(pages)


def _read_docx_text(path: Path) -> str:
    """Read paragraphs and tables from a DOCX attachment."""

    from docx import Document as DocxDocument  # type: ignore

    document = DocxDocument(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        if any(row.strip(" |") for row in rows):
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def _read_pptx_text(path: Path) -> str:
    """Read slide text and speaker notes from a PPTX attachment."""

    from pptx import Presentation  # type: ignore

    presentation = Presentation(str(path))
    slides = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        parts = [f"=== 幻灯片 {slide_number} ==="]
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                parts.append(text.strip())
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes_text = ""
        if notes_text:
            parts.append(f"备注: {notes_text}")
        if len(parts) > 1:
            slides.append("\n".join(parts))
    return "\n\n".join(slides)


_NATIVE_TEXT_READERS: dict[str, Callable[[Path], str]] = {
    ".pdf": _read_pdf_text,
    ".docx": _read_docx_text,
    ".pptx": _read_pptx_text,
}


async def _read_native_document(path: Path) -> str:
    reader = _NATIVE_TEXT_READERS.get(path.suffix.lower())
    if reader is None:
        return path.read_text(encoding="utf-8", errors="replace")
    return await asyncio.to_thread(reader, path)


async def _read_with_native_reader(path: Path) -> str:
    try:
        text = await _read_native_document(path)
    except Exception as exc:
        raise DocumentTextError(
            "parse_failed",
            f"解析 {path.name} 失败：{exc}",
        ) from exc

    if not isinstance(text, str) or not text.strip():
        raise DocumentTextError(
            "empty_content",
            f"未能从 {path.name} 提取到可用文本内容。",
        )
    return text


async def _read_with_liteparse(path: Path, *, require_libreoffice: bool) -> str:
    if not LiteParseAdapter.is_available():
        raise DocumentTextError(
            "liteparse_missing",
            (
                f"解析 {path.suffix or path.name} 需要 LiteParse。"
                f" 请安装：`{LITEPARSE_INSTALL_HINT}`。"
            ),
            install_hint=LITEPARSE_INSTALL_HINT,
        )

    if require_libreoffice and not libreoffice_available():
        hint = libreoffice_install_hint()
        raise DocumentTextError(
            "libreoffice_missing",
            (
                f"解析 {path.suffix} 需要 LibreOffice（LiteParse 内部用 soffice 做格式转换）。"
                f" 未检测到本机已安装。\n建议安装命令：{hint}"
            ),
            install_hint=hint,
        )

    adapter = LiteParseAdapter(config={"debug": False})
    try:
        text = await adapter.parse_to_text(path)
    except Exception as exc:
        msg = str(exc)
        if "LibreOffice is not installed" in msg or "soffice" in msg.lower():
            hint = libreoffice_install_hint()
            raise DocumentTextError(
                "libreoffice_missing",
                (
                    f"解析 {path.suffix} 需要 LibreOffice 做格式转换。\n"
                    f"建议安装命令：{hint}"
                ),
                install_hint=hint,
            ) from exc
        raise DocumentTextError(
            "parse_failed",
            f"LiteParse 解析 {path.name} 失败：{exc}",
        ) from exc

    if not isinstance(text, str) or not text.strip():
        raise DocumentTextError(
            "empty_content",
            f"LiteParse 未能从 {path.name} 提取到可用文本内容。",
        )
    return text


async def read_document_text(path: Path) -> str:
    """Extract plain text from a document using the lowest-dependency route.

    Routing:
      - plain text / markdown -> direct read
      - PDF / DOCX / PPTX -> native Python readers, LiteParse fallback if installed
      - DOC / PPT / XLS / XLSX / images -> LiteParse (+ LibreOffice when required)
    """
    resolved = Path(path).expanduser()
    ext = resolved.suffix.lower()

    if ext in PLAIN_TEXT_EXTS:
        return resolved.read_text(encoding="utf-8", errors="replace")

    if ext in LITEPARSE_REQUIRED_EXTS:
        return await _read_with_liteparse(
            resolved,
            require_libreoffice=ext in LIBREOFFICE_REQUIRED_EXTS,
        )

    if ext in NATIVE_READER_EXTS:
        try:
            return await _read_with_native_reader(resolved)
        except DocumentTextError as native_exc:
            if LiteParseAdapter.is_available():
                try:
                    return await _read_with_liteparse(
                        resolved,
                        require_libreoffice=False,
                    )
                except DocumentTextError:
                    raise native_exc from None
            raise

    # Unknown text-like formats retain the previous permissive fallback.
    return await _read_with_native_reader(resolved)


def read_document_text_sync(path: Path) -> str:
    """Synchronous wrapper for KB worker threads.

    Must not be called from a running asyncio event loop. Prefer
    ``await read_document_text(path)`` inside FastAPI / SSE handlers.
    """
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        return asyncio.run(read_document_text(path))
    raise RuntimeError(
        "read_document_text_sync() cannot be called from a running event loop; "
        "use await read_document_text(path) instead."
    )
